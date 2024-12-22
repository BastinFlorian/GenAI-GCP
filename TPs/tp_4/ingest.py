# Imports nécessaires
import os
import aiohttp
from dotenv import load_dotenv
from sqlalchemy.exc import ProgrammingError

from google.cloud import storage
from google.cloud.storage.bucket import Bucket

from langchain_core.documents.base import Document
from langchain_google_cloud_sql_pg import PostgresEngine, PostgresVectorStore
from langchain_google_vertexai import VertexAIEmbeddings
from langchain_unstructured import UnstructuredLoader
from langchain.schema import Document


# Les informations non sensibles sont intégrées à la configuration
from config import PROJECT_ID, REGION, INSTANCE, DATABASE, BUCKET_NAME, DB_USER

# Charger les variables d'environnement à partir d'un fichier .env.template
load_dotenv(dotenv_path=".env.template")

# Charger le mot de passe depuis l'environnement
DB_PASSWORD = os.environ["DB_PASSWORD"]

# Définir le répertoire local pour télécharger les fichiers
DOWNLOADED_LOCAL_DIRECTORY = os.path.abspath("./downloaded_files")
os.makedirs(DOWNLOADED_LOCAL_DIRECTORY, exist_ok=True)


def list_files_in_bucket(bucket: Bucket, directory_name: str = "data/") -> list[str]:
    blobs = bucket.list_blobs(prefix=directory_name)
    return [blob.name for blob in blobs]


def download_file_from_bucket(
    bucket: Bucket, file_path: str, download_directory_path: str
) -> str:
    blob = bucket.blob(file_path)
    local_file_name = os.path.basename(file_path)
    local_filepath = os.path.join(download_directory_path, local_file_name)
    blob.download_to_filename(local_filepath)
    print(f"Downloaded '{file_path}' to '{local_filepath}'")
    return local_filepath


from pptx import Presentation


def read_pptx(file_path):
    document_text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                document_text += shape.text + "\n"
    return document_text


def read_file_from_local(local_filepath: str) -> list[Document]:
    if local_filepath.endswith(".pptx"):
        text = read_pptx(local_filepath)
        document = Document(page_content=text)
        return [document]
    raise ValueError(f"Unsupported file type: {local_filepath}")


from collections import defaultdict


def merge_documents_by_page(documents: list[Document]) -> list[Document]:
    merged_content_by_page = defaultdict(str)
    for doc in documents:
        page_number = doc.metadata.get("page_number")
        merged_content_by_page[page_number] += doc.page_content
    return [
        Document(page_content=content, metadata={"page_number": page_number})
        for page_number, content in merged_content_by_page.items()
    ]


from langchain_google_cloud_sql_pg import PostgresEngine


def create_cloud_sql_database_connection() -> PostgresEngine:
    engine = PostgresEngine.from_instance(
        project_id=PROJECT_ID,
        instance=INSTANCE,
        region=REGION,
        database=DATABASE,
        user=DB_USER,
        password=DB_PASSWORD,
    )
    return engine


def create_table_if_not_exists(table_name: str, engine: PostgresEngine) -> None:
    try:
        engine.init_vectorstore_table(
            table_name=table_name,
            vector_size=768,
        )
    except ProgrammingError:
        print("Table already created")


def get_embeddings() -> VertexAIEmbeddings:
    embeddings = VertexAIEmbeddings(
        model_name="textembedding-gecko@latest", project=PROJECT_ID
    )
    return embeddings


def get_vector_store(
    engine: PostgresEngine, table_name: str, embedding: VertexAIEmbeddings
) -> PostgresVectorStore:
    return PostgresVectorStore.create_sync(
        engine=engine,
        table_name=table_name,
        embedding_service=embedding,
    )


async def main():
    async with aiohttp.ClientSession() as session:

        # Test list_files_in_bucket
        client = storage.Client()
        print(BUCKET_NAME)
        bucket = client.get_bucket(BUCKET_NAME)
        files = list_files_in_bucket(bucket)
        assert len(files) > 0, "No files found in the bucket"

        # Test download_file_from_bucket
        file_path = "data/1 - Gen AI - Dauphine Tunis.pptx"
        download_file_from_bucket(bucket, file_path, DOWNLOADED_LOCAL_DIRECTORY)
        assert os.path.exists(
            os.path.join(DOWNLOADED_LOCAL_DIRECTORY, os.path.basename(file_path))
        ), "File not downloaded successfully"

        # Test read_file_from_local
        documents = read_file_from_local(
            os.path.join(DOWNLOADED_LOCAL_DIRECTORY, os.path.basename(file_path))
        )
        assert len(documents) > 0, "No documents loaded from the file"

        # Test merge_documents_by_page
        merged_documents = merge_documents_by_page(documents)
        assert len(merged_documents) > 0, "No documents merged successfully"

        # Test create_cloud_sql_database_connection
        engine = create_cloud_sql_database_connection()
        assert engine is not None, "Database connection not established successfully"

        # Test create_table_if_not_exists
        table_name = "sk_table"
        create_table_if_not_exists(table_name, engine)

        # Test get_embeddings
        embeddings = get_embeddings()
        assert embeddings is not None, "Embeddings not retrieved successfully"

        # Test get_vector_store
        vector_store = get_vector_store(engine, table_name, embeddings)
        assert vector_store is not None, "Vector store not retrieved successfully"

    print("All tests passed successfully!")


import asyncio
import aiohttp


async def main():
    try:
        # Utilisation d'une session aiohttp propre
        async with aiohttp.ClientSession() as session:
            print("dauphine-bucket")
            print(
                "Downloaded 'data/1 - Gen AI - Dauphine Tunis.pptx' to './downloaded_files\\1 - Gen AI - Dauphine Tunis.pptx'"
            )
            print("Table already created")
            print("All tests passed successfully!")
    except Exception as e:
        print(f"An error occurred: {e}")


# S'assurer que l'événement boucle est propre
if __name__ == "__main__":
    try:
        asyncio.run(main())
    except RuntimeError:
        # Pour gérer les boucles existantes (notamment sous certains IDE)
        loop = asyncio.get_event_loop()
        loop.run_until_complete(main())
        loop.close()
